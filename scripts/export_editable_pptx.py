#!/usr/bin/env python3
"""可编辑 PPTX 导出：底层干净背景图 + 色块独立矩形 + 文字独立文本框（三层）。

设计（三层可编辑取向）：
- 第 1 层 背景图：每页整页 PNG 截图，但**去掉所有文字与实色色块**（仅保留
  图片 / 网格底纹 / 页脚 chrome / 渐变装饰 / 幻灯片底色），背景因此「更干净」。
- 第 2 层 色块矩形：从渲染后的 DOM 提取所有「实色块 / 仅边框块」
  （坐标 + 填充色 + 边框色/宽 + 圆角），在对应位置放 PPTX 矩形（autoshape），
  可单独改色 / 改边 / 移动 / 删除。
- 第 3 层 文字文本框：从 DOM 提取可见文本（坐标 + 字号 + 颜色 + 字重 + 字体 + 对齐），
  在对应位置放可编辑文本框（无填充），叠在背景与色块之上，用户在 PowerPoint
  里直接点文本框改字即可二次修改。

关键实现：
- 干净背景图用「注入 CSS 把 active slide 的所有后代 visibility:hidden，仅 IMG/SVG/
  CANVAS/.grid-lines/.chrome 及其子树保持可见」实现；用 visibility 而非 display，
  避免 flex 子项 hide 后兄弟重排、坐标错位。
- 含 <img>/<svg>/<canvas> 或 background-image（渐变/图片）的元素不参与色块提取，
  永远留在背景图里，绝不隐藏，避免丢图。
- 色块矩形绘制顺序 = DOM 文档顺序，保证嵌套块（白卡内红卡）层级正确；
  文字统一在所有色块之后添加，确保文字永远在最上层。
"""
import sys, os, re, tempfile, time, pathlib

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
if not os.path.exists(CHROME):
    CHROME = None


def count_slides(src):
    try:
        return open(src, encoding="utf-8").read().count('class="slide"')
    except Exception:
        return 0


def _font_family(ff):
    ff = (ff or "").lower()
    if "noto serif" in ff or "song" in ff:
        return "宋体"
    if "yahei" in ff or "pingfang" in ff or "hei" in ff or "noto sans" in ff:
        return "微软雅黑"
    if "playfair" in ff or "times" in ff or "georgia" in ff:
        return "Times New Roman"
    if "mono" in ff:
        return "Consolas"
    return "微软雅黑"


def _rgb(color):
    m = re.search(r"rgba?\((\d+),\s*(\d+),\s*(\d+)", color or "")
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return RGBColor(0, 0, 0)


# ---------------------------------------------------------------------------
# JS：提取可见文本块（坐标 + 字号 + 颜色 + 字重 + 字体 + 对齐），按真实视觉行拆分
# ---------------------------------------------------------------------------
def _extract_js():
    return """(function(){
      var slide=document.querySelector('section.slide.active')||document.querySelector('section.slide');
      if(!slide) return [];
      var walker=document.createTreeWalker(slide,NodeFilter.SHOW_TEXT,null);
      var groups=new Map(); var node;
      while(node=walker.nextNode()){
        var text=node.textContent; if(!text||!text.trim()) continue;
        var el=node.parentElement; if(!el) continue;
        var cs=getComputedStyle(el);
        if(cs.display==='none'||cs.visibility==='hidden'||parseFloat(cs.opacity||'1')<=0.01) continue;
        var er=el.getBoundingClientRect(); if(er.width<1||er.height<1) continue;
        if(!groups.has(el)) groups.set(el,{cs:cs,lines:[]});
        var g=groups.get(el);
        for(var i=0;i<text.length;i++){
          var ch=text[i]; if(ch==='\\n'||ch==='\\r') continue;
          var range=document.createRange();
          range.setStart(node,i); range.setEnd(node,i+1);
          var r=range.getBoundingClientRect();
          if(r.width<0.01||r.height<0.01) continue;
          var line=null;
          for(var j=0;j<g.lines.length;j++){
            if(Math.abs(g.lines[j].top-r.top)<=2.5){ line=g.lines[j]; break; }
          }
          if(!line){
            line={text:'',left:r.left,top:r.top,right:r.right,bottom:r.bottom};
            g.lines.push(line);
          }
          line.text+=ch;
          line.left=Math.min(line.left,r.left); line.top=Math.min(line.top,r.top);
          line.right=Math.max(line.right,r.right); line.bottom=Math.max(line.bottom,r.bottom);
        }
      }
      var out=[];
      groups.forEach(function(g){
        g.lines.sort(function(a,b){return a.top-b.top||a.left-b.left;});
        g.lines.forEach(function(line){
          var txt=line.text.replace(/\\s+/g,' ').trim(); if(!txt) return;
          out.push({text:txt,x:line.left,y:line.top,w:line.right-line.left,h:line.bottom-line.top,
            fs:parseFloat(g.cs.fontSize)||16,color:g.cs.color,fw:String(g.cs.fontWeight),
            ff:g.cs.fontFamily,ta:g.cs.textAlign});
        });
      });
      return out;
    })()"""


# ---------------------------------------------------------------------------
# JS：提取色块（实色填充块 / 仅边框块），按 DOM 文档顺序返回，可序列化
#   - 排除根底/网格/页脚/进度条等
#   - 含 img/svg/canvas 或 background-image(渐变/图片) 的块不提取（留在背景图）
#   - 返回 {x,y,w,h,fill:{r,g,b}|null,border:{color:{r,g,b},width:px}|null,radius:px}
# ---------------------------------------------------------------------------
def _extract_blocks_js():
    return """(function(){
      var slide=document.querySelector('section.slide.active')||document.querySelector('section.slide');
      if(!slide) return [];
      var EX={'deck-viewport':1,'deck-stage':1,slide:1,'grid-lines':1,progress:1,'progress-fill':1,'edit-toggle':1,chrome:1};
      function clsOf(el){if(!el||!el.className)return '';if(typeof el.className==='object'&&el.className.baseVal!==undefined)return el.className.baseVal;return ''+el.className;}
      function parseRGB(c){if(!c)return null;var s=String(c).trim();
        // 兼容 'rgb(255, 51, 0)' / 'rgb(255 51 0)' / 'rgba(255, 51, 0, 1)'
        var m=s.match(/^rgba?\\(\\s*([\\d.]+)\\s*[, ]\\s*([\\d.]+)\\s*[, ]\\s*([\\d.]+)/);
        return m?{r:+m[1],g:+m[2],b:+m[3]}:null;}
      function alpha0(c){if(!c)return true;var s=String(c).trim();
        // 只有 rgba(...) 才可能透明;rgb(...) 没有 alpha 通道
        if(!/^rgba\\(/.test(s))return false;
        var p=s.slice(5,-1).split(',');
        return p.length>=4&&parseFloat(p[p.length-1])<=0.001;}
      var out=[];
      var w=document.createTreeWalker(slide,NodeFilter.SHOW_ELEMENT,null);
      var el;
      while(el=w.nextNode()){
        var toks=(' '+clsOf(el)+' ').split(/\\s+/);
        for(var t=0;t<toks.length;t++){ if(EX[toks[t]]){ el=null; break; } }
        if(!el) continue;
        if(el.tagName==='HTML'||el.tagName==='BODY') continue;
        var cs=getComputedStyle(el);
        var er=el.getBoundingClientRect();
        if(er.width<6||er.height<6) continue;
        if(el.querySelector('img,svg,canvas')) continue;            // 含媒体 => 留在背景
        if(cs.backgroundImage && cs.backgroundImage!=='none') continue; // 渐变/图片背景 => 留在背景
        var fill=parseRGB(cs.backgroundColor);
        var hasFill = fill && !alpha0(cs.backgroundColor);
        var bw=parseFloat(cs.borderTopWidth)||0;
        var bcol=parseRGB(cs.borderTopColor);
        var hasBorder = bw>0 && bcol && !alpha0(cs.borderTopColor);
        if(!hasFill && !hasBorder) continue;
        var br=parseFloat(cs.borderTopLeftRadius)||0;
        out.push({x:er.left,y:er.top,w:er.width,h:er.height,
          fill: hasFill?fill:null,
          border: hasBorder?{color:bcol,width:bw}:null,
          radius: br});
      }
      return out;
    })()"""


# ---------------------------------------------------------------------------
# JS：注入「干净背景」隐藏样式（隐藏 active slide 全部后代，仅保留媒体/网格/页脚）
# ---------------------------------------------------------------------------
def _inject_clean_js():
    return """(function(){
      var slide=document.querySelector('section.slide.active')||document.querySelector('section.slide');
      if(!slide) return;
      slide.setAttribute('data-clean','1');
      var s=document.getElementById('cleanhide'); if(s) s.remove();
      s=document.createElement('style'); s.id='cleanhide';
      s.textContent="section.slide[data-clean] :not(IMG, SVG, CANVAS, .grid-lines, .chrome, .chrome *):not(:has(IMG, SVG, CANVAS)),"
        +".progress, .progress-fill, .edit-toggle"
        +"{visibility:hidden !important;}";
      document.head.appendChild(s);
    })()"""


def _remove_clean_js():
    return """(function(){
      var s=document.getElementById('cleanhide'); if(s) s.remove();
      var sl=document.querySelector('section.slide[data-clean]');
      if(sl) sl.removeAttribute('data-clean');
    })()"""


def build_editable_pptx(src):
    """渲染每页 + 提取文本/色块 + 生成三层可编辑 PPTX，返回输出路径。"""
    n = count_slides(src)
    if n == 0:
        return None
    src = pathlib.Path(src)
    WORK = src.parent
    # 渲染 + 提取（复用 deck.showSlide + 截图 + DOM 文本/色块提取）
    html = open(src, encoding="utf-8").read()
    html = re.sub(r'<link[^>]*fonts\.(googleapis|gstatic)\.com[^>]*>', "", html)
    tmpf = tempfile.NamedTemporaryFile(prefix="fde-deck-", suffix=".html",
                                       dir=str(WORK), delete=False)
    TMPSRC = pathlib.Path(tmpf.name)
    tmpf.close()
    TMPSRC.write_text(html, encoding="utf-8")
    args = [
        "--no-sandbox", "--no-first-run", "--no-default-browser-check",
        "--disable-dev-shm-usage", "--disable-gpu", "--disable-component-update",
        "--disable-background-networking", "--disable-sync", "--no-pings",
        "--disable-features=Translate,BackForwardCache,OptimizationHints,MediaRouter,InfiniteSessionRestore",
    ]
    data = []
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser = p.chromium.launch(executable_path=CHROME, headless=True, args=args)
            page = browser.new_page(viewport={"width": 1920, "height": 1080},
                                    device_scale_factor=1)
            page.goto(f"file://{TMPSRC}", wait_until="domcontentloaded")
            page.wait_for_function("typeof deck !== 'undefined'", timeout=15000)
            for i in range(n):
                page.evaluate(f"deck.showSlide({i})")
                page.wait_for_timeout(800)
                # 1) 完整预览图（供编辑器缩略图/预览，含文字+色块，不变）
                full = WORK / f"preview-P{i+1:02d}.png"
                page.screenshot(path=str(full),
                                clip={"x": 0, "y": 0, "width": 1920, "height": 1080})
                # 2) 文本块（正常可见状态提取，坐标基于真实布局）
                text_blocks = page.evaluate(_extract_js())
                # 3) 色块矩形（正常可见状态提取）
                block_list = page.evaluate(_extract_blocks_js())
                # 4) 干净背景图：隐藏文字+实色块，仅留图片/网格/页脚 => visibility:hidden
                page.evaluate(_inject_clean_js())
                page.wait_for_timeout(120)
                bg = WORK / f"_bg_P{i+1:02d}.png"
                page.screenshot(path=str(bg),
                                clip={"x": 0, "y": 0, "width": 1920, "height": 1080})
                page.evaluate(_remove_clean_js())
                data.append((str(full), str(bg), text_blocks, block_list))
            browser.close()
    finally:
        try:
            TMPSRC.unlink()
        except Exception:
            pass

    # 构建三层可编辑 PPTX
    prs = Presentation()
    # 锁定 16:9 / 1920×1080（标准宽屏）。EMU: 1 inch = 914400 EMU，
    # 16:9 = 13.333in × 7.5in = 12192000 × 6858000 EMU（精确整数，
    # 不用 Inches(13.333) 的浮点截断，避免日后被误改成 4:3）。
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    SX = 12192000 / 1920.0 / 914400.0   # 英寸/px（横向）
    SY = 6858000 / 1080.0 / 914400.0    # 英寸/px（纵向）
    PT_PER_PX = SX * 72.0         # 点/px：0.5 pt/px（与位置同源，避免字被放大）

    for full, bg, text_blocks, block_list in data:
        slide = prs.slides.add_slide(blank)
        # 第 1 层：干净背景图
        slide.shapes.add_picture(bg, 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        # 第 2 层：色块独立矩形（DOM 顺序保证嵌套层级正确）
        for b in block_list:
            left = Inches(b["x"] * SX)
            top = Inches(b["y"] * SY)
            width = Inches(max(b["w"] * SX, 0.05))
            height = Inches(max(b["h"] * SY, 0.05))
            # 防止溢出页面
            if left + width > prs.slide_width:
                width = prs.slide_width - left
            if top + height > prs.slide_height:
                height = prs.slide_height - top
            is_rounded = (b.get("radius") or 0) > 4
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if is_rounded else MSO_SHAPE.RECTANGLE
            shp = slide.shapes.add_shape(shape_type, left, top, width, height)
            # 填充
            if b["fill"]:
                f = b["fill"]
                shp.fill.solid()
                shp.fill.fore_color.rgb = RGBColor(f["r"], f["g"], f["b"])
            else:
                shp.fill.background()   # 无填充（仅边框块，如 .duo .panel）
            # 边框
            if b["border"]:
                bd = b["border"]
                shp.line.color.rgb = RGBColor(bd["color"]["r"], bd["color"]["g"], bd["color"]["b"])
                shp.line.width = Pt(bd["width"] * PT_PER_PX)
            else:
                shp.line.fill.background()   # 无边框
            # 文本框垂直居中（可选，保持顶部对齐更贴近原版）
            shp.text_frame.word_wrap = False
        # 第 3 层：文字独立文本框（无填充，永远在最上层）
        for b in text_blocks:
            left = Inches(b["x"] * SX)
            top = Inches(b["y"] * SY)
            max_width_px = max(1920.0 - b["x"], 1.0)
            safe_width_px = min(b["w"] * 1.08 + 4.0, max_width_px)
            width = Inches(max(safe_width_px * SX, 0.15))
            height = Inches(max(b["h"] * SY * 1.10, 0.15))
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            tf.word_wrap = False   # 每个框就是浏览器渲染后的一整行，禁止 PowerPoint 二次换行
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP
            para = tf.paragraphs[0]
            para.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }.get(b["ta"], PP_ALIGN.LEFT)
            run = para.add_run()
            run.text = b["text"]
            run.font.size = Pt(b["fs"] * PT_PER_PX)
            run.font.bold = (b["fw"] in ("bold", "700", "800", "900") or
                             (b["fw"].isdigit() and int(b["fw"]) >= 700))
            run.font.color.rgb = _rgb(b["color"])
            run.font.name = _font_family(b["ff"])
    out = WORK / (src.stem + ".pptx")
    prs.save(str(out))
    return str(out)


def build_image_pptx(src, dsf=2):
    """高清整页图片版：每页整页高分辨率截图直接铺满一页，无文本框（不可二次改字）。

    适合「只要高清图、拿去演示 / 打印、不打算在 PowerPoint 里改字」的场景。
    dsf 默认 2 → 1920×1080 视口 ×2 = 3840×2160 实际像素，文字 / 配图更锐利。
    """
    n = count_slides(src)
    if n == 0:
        return None
    src = pathlib.Path(src)
    WORK = src.parent
    html = open(src, encoding="utf-8").read()
    html = re.sub(r'<link[^>]*fonts\.(googleapis|gstatic)\.com[^>]*>', "", html)
    tmpf = tempfile.NamedTemporaryFile(prefix="fde-deck-", suffix=".html",
                                       dir=str(WORK), delete=False)
    TMPSRC = pathlib.Path(tmpf.name)
    tmpf.close()
    TMPSRC.write_text(html, encoding="utf-8")
    args = [
        "--no-sandbox", "--no-first-run", "--no-default-browser-check",
        "--disable-dev-shm-usage", "--disable-gpu", "--disable-component-update",
        "--disable-background-networking", "--disable-sync", "--no-pings",
        "--disable-features=Translate,BackForwardCache,OptimizationHints,MediaRouter,InfiniteSessionRestore",
    ]
    shots = []
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser = p.chromium.launch(executable_path=CHROME, headless=True, args=args)
            page = browser.new_page(viewport={"width": 1920, "height": 1080},
                                    device_scale_factor=dsf)
            page.goto(f"file://{TMPSRC}", wait_until="domcontentloaded")
            page.wait_for_function("typeof deck !== 'undefined'", timeout=15000)
            for i in range(n):
                page.evaluate(f"deck.showSlide({i})")
                page.wait_for_timeout(800)
                shot = WORK / f"preview-P{i+1:02d}.png"
                page.screenshot(path=str(shot),
                                clip={"x": 0, "y": 0, "width": 1920, "height": 1080})
                shots.append(shot)
            browser.close()
    finally:
        try:
            TMPSRC.unlink()
        except Exception:
            pass

    prs = Presentation()
    # 锁定 16:9 / 1920×1080（标准宽屏）。与可编辑版同款精确 EMU，
    # 保证「高清整页图」导出也是 16:9，绝不退化成 4:3。
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    for shot in shots:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(shot), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
    out = WORK / (src.stem + ".pptx")
    prs.save(str(out))
    return str(out)


if __name__ == "__main__":
    if len(sys.argv) > 1:
        print(build_editable_pptx(sys.argv[1]))
