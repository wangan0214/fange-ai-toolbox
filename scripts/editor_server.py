#!/usr/bin/env python3
"""FDE 幻灯片本地编辑器 —— 后端服务。

功能：
  GET  /                   -> 编辑器页面
  GET  /api/files          -> 列出可编辑的 HTML 文件（相对 ALLOWED_ROOT）
  GET  /api/load?file=     -> 读取指定源 HTML 文件原文
  POST /api/save           -> 把前端传回的 HTML 原子写回源文件（并自动存历史快照）
  POST /api/render         -> 后台渲染当前文件全部 PNG（单会话 Playwright）
  POST /api/export-pptx    -> 后台渲染 + 打包成 .pptx（16:9）
  GET  /api/download-pptx  -> 下载生成的 .pptx
  GET  /api/render-status  -> 渲染/导出进度
  GET  /api/root           -> 返回当前实例的 root（供一键启动脚本检测并自动重启）
  GET  /api/history?file=  -> 列出该文件的历史快照
  GET  /api/history-file?file=&snap= -> 读取某个历史快照内容（预览/下载）
  POST /api/rollback       -> 把某个历史快照原子写回源文件（回滚前自动留快照）
  POST /api/upload-feishu  -> 把当前 HTML 或已导出 PPTX 上传到飞书云空间（lark-cli drive +upload）

纯本地 (127.0.0.1)，数据零外发；仅「上传飞书云空间」会调用本机 lark-cli 与飞书交互。
"""
import http.server
import socketserver
import json
import os
import sys
import subprocess
import threading
import re
import time
import argparse
from urllib.parse import quote, urlparse, parse_qs
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
try:
    from export_editable_pptx import build_editable_pptx, build_image_pptx
except Exception:
    build_editable_pptx = None
    build_image_pptx = None

# ---------- 默认值（可被 CLI 参数覆盖） ----------
# 编辑器自家目录：scripts/ 的上级（即「帆哥 PPT 编辑器/」）。
# 启动后 ALLOWED_ROOT 会随「打开的 deck」动态切到该 deck 所在目录（原地编辑），
# 但历史快照 / 渲染产物 / 日志一律存编辑器自家目录，绝不污染用户项目文件夹。
EDITOR_HOME = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_DEFAULT_ROOT = EDITOR_HOME
_DEFAULT_PORT = 8731

def _resolve_default_file(root):
    """根目录下第一个 .html 文件；都没有则用 'index.html' 占位（用户可在 UI 里改）。"""
    try:
        for fn in sorted(os.listdir(root)):
            if fn.endswith(".html") and not fn.startswith("fde-deck-"):
                return fn
    except OSError:
        pass
    return "index.html"

def _parse_args():
    p = argparse.ArgumentParser(
        description="FDE 幻灯片本地编辑器（HTML 幻灯片可视化编辑器 + 单会话渲染 + 历史快照回滚 + PPT 导出）",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="示例：python editor_server.py --root ~/decks --port 9000"
    )
    p.add_argument("--root", default=_DEFAULT_ROOT,
                   help=f"HTML 文件所在根目录（默认：{_DEFAULT_ROOT}）")
    p.add_argument("--port", type=int, default=_DEFAULT_PORT,
                   help=f"HTTP 端口（默认：{_DEFAULT_PORT}）")
    p.add_argument("--default", dest="default_file", default=None,
                   help="默认打开的文件名（不传则用根目录下第一个 .html）")
    return p.parse_args()

_args = _parse_args()
ALLOWED_ROOT = os.path.abspath(_args.root)   # 初始 = 编辑器自家目录；打开 deck 后切到该 deck 目录
PORT = _args.port
DEFAULT_FILE = _args.default_file or "index.html"
RENDER_PW = os.path.join(EDITOR_HOME, "scripts", "render_slides_pw.py")
HISTORY_BASE = os.path.join(EDITOR_HOME, ".history")   # 历史快照存这里，不进用户文件夹
RENDER_BASE = os.path.join(EDITOR_HOME, ".render")     # 渲染/导出产物存这里，不进用户文件夹
ACTIVE_FILE = None     # 当前打开的 deck 绝对路径（None = 还没打开）
ACTIVE_DIR = None      # 当前打开的 deck 所在目录
ROOT = ALLOWED_ROOT                # 静态文件根（未打开时是编辑器自家目录；打开后是 deck 目录，用于服务 deck 配图）
# 编辑器前端页面（editor.html）随 skill 走，固定在本文件同级 templates/，不依赖用户 decks 目录
SKILL_SCRIPTS = os.path.dirname(os.path.abspath(__file__))   # scripts/
EDITOR_HTML = os.path.join(SKILL_SCRIPTS, "..", "templates", "editor.html")

# 渲染状态（进程内共享；字典里不放锁实例，避免 JSON 序列化失败）
render_status = {"running": False, "total": 0, "done": 0, "error": "", "last": "",
                 "cancelled": False, "phase": "", "pptx_ready": False, "pptx_name": ""}
render_lock = threading.Lock()      # 单独模块级锁，保护 render_status 和 proc 的并发读写

MAX_SNAPSHOTS = 60                  # 每个文件保留的历史快照上限（超出删最旧，单文件删除）


# ---------- 路径安全 ----------
def safe_abs(rel):
    """校验相对路径在 ALLOWED_ROOT 内且为 .html，返回绝对路径或 None。"""
    if not rel:
        return None
    if ".." in rel.replace("\\", "/").split("/"):
        return None
    if not rel.endswith(".html"):
        return None
    fp = os.path.normpath(os.path.join(ALLOWED_ROOT, rel))
    base = os.path.normpath(ALLOWED_ROOT)
    if fp != base and not fp.startswith(base + os.sep):
        return None
    return fp


def base_of(rel):
    return os.path.splitext(os.path.basename(rel))[0]


def safe_snap(snap):
    """快照名只允许文件名（禁路径分隔与 ..），允许中文等非 ASCII 字符。"""
    if not snap or not snap.endswith(".html"):
        return False
    if "/" in snap or "\\" in snap or ".." in snap:
        return False
    return True


def resolve_file(rel):
    """原地编辑模式下的「目标源文件」解析：打开 deck 后，所有读写一律落在 ACTIVE_FILE
    所在目录（deck 目录）内，绝不逃逸到编辑器目录或用户其他位置。

    - 未打开 deck（ACTIVE_FILE 为 None）：一律返回 None（前端引导页不会发起写操作）。
    - rel 为空或等同 DEFAULT_FILE：返回 ACTIVE_FILE 本身。
    - rel 为绝对路径：必须是 ACTIVE_FILE 本身，或同一 deck 目录下的其他 .html（切换文件）。
    - rel 为文件名：必须是同一 deck 目录下的 .html。
    任何越界都返回 None（handler 据此返回 400）。
    """
    if not ACTIVE_FILE:
        return None
    if not rel or rel == DEFAULT_FILE:
        return ACTIVE_FILE
    if os.path.isabs(rel):
        if rel == ACTIVE_FILE:
            return ACTIVE_FILE
        if os.path.dirname(rel) == ACTIVE_DIR and rel.endswith(".html"):
            return rel
        return None
    cand = os.path.normpath(os.path.join(ALLOWED_ROOT, rel))
    if os.path.dirname(cand) == ACTIVE_DIR and cand.endswith(".html"):
        return cand
    return None


# ---------- 历史快照（存编辑器自家目录，按 deck 绝对路径命名空间）----------
def deck_key(abs_path):
    """用 deck 绝对路径生成稳定且唯一的子目录名，避免不同目录下同名 index.html 互相覆盖。"""
    import hashlib
    h = hashlib.sha1(abs_path.encode("utf-8")).hexdigest()[:12]
    return f"{h}_{base_of(abs_path)}"


def snapshot_dir():
    d = os.path.join(HISTORY_BASE, deck_key(ACTIVE_FILE))
    os.makedirs(d, exist_ok=True)
    return d


def save_snapshot(rel, html):
    """把 html 存一份带时间戳的快照；超出上限删最旧（一次一个明确路径，合规）。

    时间戳精确到微秒，避免同秒内连续保存/回滚时文件名碰撞导致快照被覆盖。
    快照存于 HISTORY_BASE/<deckKey>/，绝不写入用户项目文件夹。
    """
    d = snapshot_dir()
    ts = time.strftime("%Y%m%dT%H%M%S") + ("%06d" % int((time.time() % 1) * 1e6))
    name = f"{base_of(ACTIVE_FILE)}.{ts}.html"
    with open(os.path.join(d, name), "w", encoding="utf-8") as f:
        f.write(html)
    prefix = base_of(ACTIVE_FILE) + "."
    files = sorted(
        [os.path.join(d, f) for f in os.listdir(d)
         if f.endswith(".html") and f.startswith(prefix)],
        key=os.path.getmtime)
    while len(files) > MAX_SNAPSHOTS:
        old = files.pop(0)
        try:
            os.remove(old)
        except Exception:
            pass


def list_history():
    d = os.path.join(HISTORY_BASE, deck_key(ACTIVE_FILE))
    if not os.path.isdir(d):
        return []
    prefix = base_of(ACTIVE_FILE) + "."
    items = []
    for f in os.listdir(d):
        if f.endswith(".html") and f.startswith(prefix):
            fp = os.path.join(d, f)
            items.append({"name": f, "mtime": os.path.getmtime(fp),
                          "size": os.path.getsize(fp)})
    items.sort(key=lambda x: x["mtime"], reverse=True)
    return items


def list_html_files():
    """未打开 deck 时返回空（编辑器不再罗列某个固定根下的文件）；
    打开 deck 后罗列当前 deck 目录下所有 .html（含子目录），供「切换文件」下拉使用。"""
    res = []
    for dp, dirs, fns in os.walk(ALLOWED_ROOT):
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for fn in fns:
            if fn.endswith(".html") and not fn.startswith("fde-deck-"):
                full = os.path.join(dp, fn)
                rel = os.path.relpath(full, ALLOWED_ROOT)
                res.append({"rel": rel, "name": fn, "abs": full})
    res.sort(key=lambda x: (x["rel"] != DEFAULT_FILE, x["rel"]))
    return res


# ---------- 打开 deck（原地编辑核心）----------
RECENT_FILE = os.path.join(EDITOR_HOME, ".recent.json")

def read_recent():
    try:
        with open(RECENT_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, list):
            return data
    except Exception:
        pass
    return []

def push_recent(abs_path):
    lst = read_recent()
    lst = [p for p in lst if p != abs_path]
    lst.insert(0, abs_path)
    lst = lst[:12]
    try:
        with open(RECENT_FILE, "w", encoding="utf-8") as f:
            json.dump(lst, f, ensure_ascii=False, indent=2)
    except Exception:
        pass
    return lst

def remove_recent(abs_path):
    lst = read_recent()
    lst = [p for p in lst if p != abs_path]
    try:
        with open(RECENT_FILE, "w", encoding="utf-8") as f:
            json.dump(lst, f, ensure_ascii=False, indent=2)
    except Exception:
        pass
    return lst

def set_active(abs_path):
    """打开一个 deck：把 ALLOWED_ROOT 动态切到该文件所在目录（原地编辑）。
    之后所有现有 handler（load/save/render/rollback）都作用于这个 deck，无需复制文件。"""
    global ALLOWED_ROOT, DEFAULT_FILE, ACTIVE_FILE, ACTIVE_DIR, ROOT
    abs_path = os.path.abspath(abs_path)
    ALLOWED_ROOT = os.path.dirname(abs_path)
    DEFAULT_FILE = os.path.basename(abs_path)
    ACTIVE_FILE = abs_path
    ACTIVE_DIR = os.path.dirname(abs_path)
    ROOT = os.path.dirname(abs_path)   # 静态资源（配图等）也以 deck 目录为根，保证 relative 引用可加载

def pick_path(kind):
    """调用 macOS 原生选择对话框，返回所选文件/文件夹的绝对路径；取消或失败返回 None。"""
    if kind == "folder":
        script = ('try\n'
                  'POSIX path of (choose folder with prompt "选择要编辑的 deck 文件夹")\n'
                  'on error\nreturn ""\nend try')
    else:
        script = ('try\n'
                  'POSIX path of (choose file with prompt "选择要编辑的 PPT（index.html 等）")\n'
                  'on error\nreturn ""\nend try')
    try:
        out = subprocess.run(["osascript", "-e", script],
                             capture_output=True, text=True, timeout=120).stdout.strip()
    except Exception:
        return None
    if not out or out.lower().startswith("null"):
        return None
    return out.rstrip("/") or None


# ---------- 渲染 ----------
def count_slides(file):
    try:
        return open(file, "r", encoding="utf-8").read().count('class="slide"')
    except Exception:
        return 0


def render_all_pw(file):
    """单会话渲染：用 Playwright 启动一次 Chrome，把全部幻灯片截成 PNG。"""
    n = count_slides(file)
    outdir = os.path.join(RENDER_BASE, deck_key(file))   # 渲染产物存编辑器自家目录
    os.makedirs(outdir, exist_ok=True)
    with render_lock:
        render_status["total"] = n
        render_status["done"] = 0
        render_status["cancelled"] = False
    proc = subprocess.Popen(
        [sys.executable, RENDER_PW, "--out", outdir, file, "1", str(n)],
        cwd=ACTIVE_DIR or os.path.dirname(file),
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
        bufsize=1,
    )
    with render_lock:
        render_status["_proc"] = proc
    try:
        for line in proc.stdout:
            line = line.strip()
            m = re.match(r"^P(\d+) ->", line)
            if m:
                with render_lock:
                    render_status["done"] = int(m.group(1))
                    render_status["last"] = line
            if render_status["cancelled"]:
                try:
                    proc.kill()
                except Exception:
                    pass
                break
        try:
            proc.wait(timeout=15)
        except subprocess.TimeoutExpired:
            try:
                proc.kill()
            except Exception:
                pass
            try:
                proc.wait(timeout=5)
            except Exception:
                pass
    finally:
        with render_lock:
            render_status["_proc"] = None
    return not render_status["cancelled"]


def build_pptx(file):
    """把 preview-P01..P{n:02d}.png 打包成 16:9 的 .pptx，返回输出路径。"""
    from pptx import Presentation
    from pptx.util import Emu
    n = count_slides(file)
    outdir = os.path.join(RENDER_BASE, deck_key(file))   # 产物存编辑器自家目录
    os.makedirs(outdir, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 锁定 16:9 / 1920×1080（标准宽屏，精确 EMU）
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    for i in range(1, n + 1):
        img = os.path.join(outdir, f"preview-P{i:02d}.png")
        if not os.path.exists(img):
            continue
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img, 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
    out = os.path.join(outdir, base_of(file) + ".pptx")
    prs.save(out)
    return out


def run_export(file, mode="editable"):
    with render_lock:
        render_status["running"] = True
        render_status["error"] = ""
        render_status["cancelled"] = False
        render_status["phase"] = "building"
        render_status["last"] = ""
        render_status["pptx_ready"] = False
        render_status["pptx_name"] = ""
    try:
        if mode == "image":
            # 高清整页图片版：每页整页高分辨率截图直接铺满一页，无文本框（不可二次改字）
            if build_image_pptx is not None:
                out = build_image_pptx(file, out_dir=os.path.join(RENDER_BASE, deck_key(file)))
            else:
                render_all_pw(file)
                out = build_pptx(file)
            label = "高清整页图"
        else:
            # 可编辑版（默认）：整页截图背景 + 文字转可编辑文本框
            if build_editable_pptx is not None:
                out = build_editable_pptx(file, out_dir=os.path.join(RENDER_BASE, deck_key(file)))
            else:
                render_all_pw(file)
                out = build_pptx(file)
            label = "可编辑版"
        if out and not render_status["cancelled"]:
            with render_lock:
                render_status["phase"] = "pptx_done"
                render_status["pptx_ready"] = True
                render_status["pptx_name"] = os.path.basename(out)
                render_status["last"] = f"PPT 已生成（{label}）"
        elif render_status["cancelled"]:
            with render_lock:
                render_status["phase"] = ""
    except Exception as e:
        with render_lock:
            if not render_status["cancelled"]:
                render_status["error"] = str(e)
                render_status["phase"] = "pptx_error"
    finally:
        with render_lock:
            render_status["running"] = False


def run_upload_feishu(kind):
    """上传当前打开 deck 的 HTML 或已导出的 PPTX 到飞书云空间（Drive 根目录）。

    返回 {ok, url, token, msg, raw}。url 为可打开的飞书云空间链接。
    所有读写都基于 ACTIVE_FILE / RENDER_BASE，绝不触碰用户其他文件。
    """
    import shutil
    if kind == "pptx":
        name = render_status.get("pptx_name", "")
        if not name:
            return {"ok": False, "msg": "尚未导出 PPT，请先点「导出 PPT」", "raw": ""}
        if not ACTIVE_FILE:
            return {"ok": False, "msg": "尚未打开 deck", "raw": ""}
        fp = os.path.join(RENDER_BASE, deck_key(ACTIVE_FILE), name)
        if not os.path.isfile(fp):
            return {"ok": False, "msg": "PPT 文件不存在，请重新导出", "raw": ""}
    else:
        fp = ACTIVE_FILE
        if not fp or not os.path.isfile(fp):
            return {"ok": False, "msg": "尚未打开 deck 或源文件不存在", "raw": ""}
    cli = shutil.which("lark-cli") or "lark-cli"
    # lark-cli 要求 --file 为「当前目录下的相对路径」，故切到文件所在目录再传文件名
    up_dir = os.path.dirname(fp)
    up_name = os.path.basename(fp)
    try:
        proc = subprocess.run(
            [cli, "drive", "+upload", "--file", up_name, "--as", "user"],
            cwd=up_dir, capture_output=True, text=True, timeout=180)
    except Exception as e:
        return {"ok": False, "msg": f"调用 lark-cli 失败：{e}", "raw": ""}
    out = (proc.stdout or "") + "\n" + (proc.stderr or "")
    url = token = ""
    # 1) 优先解析 JSON（lark-cli 通常返回 {data:{url:...}} 之类）
    try:
        j = json.loads(proc.stdout.strip())
        def _dig(o):
            if isinstance(o, dict):
                for k, v in o.items():
                    if k in ("url", "url_token", "file_token", "token") and isinstance(v, str) and v:
                        return v
                    r = _dig(v)
                    if r:
                        return r
            elif isinstance(o, list):
                for it in o:
                    r = _dig(it)
                    if r:
                        return r
            return None
        found = _dig(j)
        if found:
            if re.search(r"feishu|larksuite", found):
                url = found
            else:
                token = found
    except Exception:
        pass
    # 2) 退而求其次：正则抓链接
    if not url:
        m = re.search(r"https?://[^\s\"'\\<>]+feishu\.cn[^\s\"'\\<>]*", out)
        if not m:
            m = re.search(r"https?://[^\s\"'\\<>]+larksuite\.cn[^\s\"'\\<>]*", out)
        if m:
            url = m.group(0)
    ok = (proc.returncode == 0) and (bool(url) or bool(token))
    if ok and not url and token:
        url = f"https://www.feishu.cn/drive/home/{token}"
    return {"ok": ok, "url": url, "token": token,
            "msg": ("已上传到飞书云空间" if ok else "上传失败，详见 raw"),
            "raw": out.strip()[-1600:]}


def run_feishu_auth():
    """检测本机 lark-cli 是否已授权飞书用户身份（且具备 drive:file:upload 上传权限）。

    返回 {bound, identity, hasDriveScope, msg}。
    - bound=True   → 可直接上传
    - bound=False  → 前端应弹出激活/授权引导（让用户执行 `lark-cli auth login`）
    """
    import shutil, re
    cli = shutil.which("lark-cli") or "lark-cli"
    try:
        proc = subprocess.run([cli, "auth", "status"],
                              capture_output=True, text=True, timeout=30)
    except Exception as e:
        return {"bound": False, "identity": "", "hasDriveScope": False,
                "msg": f"未检测到 lark-cli：{e}（请先安装 lark-cli）"}
    out = (proc.stdout or "") + "\n" + (proc.stderr or "")
    # 解析 JSON（容错：抓第一个 {...} 子串）
    raw = out.strip()
    j = {}
    try:
        j = json.loads(raw)
    except Exception:
        m = re.search(r"\{.*\}", raw, re.S)
        if m:
            try:
                j = json.loads(m.group(0))
            except Exception:
                j = {}
    user = (j.get("identities") or {}).get("user") or {}
    if not user:
        return {"bound": False, "identity": "",
                "hasDriveScope": False,
                "msg": "未绑定飞书用户身份（user identity 缺失），请先授权"}
    if user.get("tokenStatus", "") != "valid" or user.get("status", "") != "ready":
        return {"bound": False, "identity": user.get("userName", ""),
                "hasDriveScope": False,
                "msg": "飞书账号未处于有效授权状态，请重新授权"}
    scope = user.get("scope", "") or ""
    if "drive:file:upload" not in scope:
        return {"bound": False, "identity": user.get("userName", ""),
                "hasDriveScope": False,
                "msg": "已授权但缺少 drive:file:upload 权限，无法上传到云空间"}
    return {"bound": True, "identity": user.get("userName", ""),
            "hasDriveScope": True, "msg": "已授权飞书账号"}


def run_render(file):
    with render_lock:
        render_status["running"] = True
        render_status["error"] = ""
        render_status["cancelled"] = False
        render_status["phase"] = ""
    try:
        render_all_pw(file)
        with render_lock:
            if not render_status["cancelled"]:
                render_status["last"] = f"渲染完成 {render_status['done']} 页"
    except Exception as e:
        with render_lock:
            if not render_status["cancelled"]:
                render_status["error"] = str(e)
                render_status["phase"] = "render_error"
    finally:
        with render_lock:
            render_status["running"] = False


class Handler(http.server.BaseHTTPRequestHandler):
    def _send(self, code, body, ctype="application/json"):
        if isinstance(body, str):
            body = body.encode("utf-8")
        self.send_response(code)
        self.send_header("Content-Type", ctype)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Content-Length", str(len(body)))
        # 强制所有响应都不缓存（用户最痛的问题：硬刷都清不掉的 Chrome disk cache）
        self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
        self.send_header("Pragma", "no-cache")
        self.send_header("Expires", "0")
        self.end_headers()
        self.wfile.write(body)

    def _qs(self):
        # http.server 把 URL 路径按 Latin-1 解码，含 UTF-8 中文时会被污染。
        # 先按 latin-1 还原字节，再按 utf-8 解码，可同时兼容「裸 UTF-8」与「百分号编码」两种请求。
        try:
            raw = self.path.encode("latin-1").decode("utf-8")
        except (UnicodeEncodeError, UnicodeDecodeError):
            raw = self.path
        return parse_qs(urlparse(raw).query)

    def do_GET(self):
        if self.path in ("/", "/editor.html"):
            try:
                with open(EDITOR_HTML, "r", encoding="utf-8") as f:
                    self._send(200, f.read(), "text/html; charset=utf-8")
            except Exception as e:
                self._send(500, str(e))
        elif self.path == "/api/files":
            self._send(200, json.dumps(list_html_files(), ensure_ascii=False))
        elif self.path == "/api/recent":
            self._send(200, json.dumps(read_recent(), ensure_ascii=False))
        elif self.path.startswith("/api/pick"):
            qs = self._qs()
            kind = (qs.get("type") or ["file"])[0]
            p = pick_path(kind)
            if p is None:
                self._send(200, json.dumps({"path": None}))
            else:
                self._send(200, json.dumps({"path": p}, ensure_ascii=False))
        elif self.path.startswith("/api/open"):
            qs = self._qs()
            p = (qs.get("path") or [""])[0].rstrip("/")
            if not p:
                self._send(400, json.dumps({"error": "缺少 path"}))
                return
            # 若是目录，找里面的 deck html（index.html 优先）
            if os.path.isdir(p):
                cand = None
                for fn in ("index.html", "INDEX.HTML"):
                    if os.path.isfile(os.path.join(p, fn)):
                        cand = os.path.join(p, fn); break
                if cand is None:
                    try:
                        for fn in sorted(os.listdir(p)):
                            if fn.endswith(".html") and not fn.startswith("fde-deck-"):
                                cand = os.path.join(p, fn); break
                    except Exception:
                        pass
                p = cand
            if not p or not os.path.isfile(p) or not p.endswith(".html"):
                self._send(400, json.dumps({"error": "不是有效的 .html deck 文件"}))
                return
            try:
                with open(p, "r", encoding="utf-8") as f:
                    text = f.read()
            except Exception as e:
                self._send(500, json.dumps({"error": str(e)}))
                return
            set_active(p)
            push_recent(p)
            self._send(200, text, "text/plain; charset=utf-8")
        elif self.path.startswith("/api/load"):
            qs = self._qs()
            rel = (qs.get("file") or [DEFAULT_FILE])[0]
            fp = resolve_file(rel)
            if not fp:
                self._send(400, json.dumps({"error": "非法文件或未打开 deck"}))
                return
            try:
                with open(fp, "r", encoding="utf-8") as f:
                    self._send(200, f.read(), "text/plain; charset=utf-8")
            except Exception as e:
                self._send(500, json.dumps({"error": str(e)}))
        elif self.path.startswith("/api/history-file"):
            qs = self._qs()
            snap = (qs.get("snap") or [""])[0]
            if not ACTIVE_FILE or not safe_snap(snap):
                self._send(400, "bad request")
                return
            snap_fp = os.path.join(HISTORY_BASE, deck_key(ACTIVE_FILE), snap)
            sb = os.path.normpath(snap_fp)
            hb = os.path.normpath(HISTORY_BASE)
            if not sb.startswith(hb + os.sep) or not os.path.isfile(sb):
                self._send(404, "not found")
                return
            with open(sb, "r", encoding="utf-8") as f:
                self._send(200, f.read(), "text/html; charset=utf-8")
        elif self.path.startswith("/api/history"):
            if not ACTIVE_FILE:
                self._send(200, json.dumps([]))
                return
            self._send(200, json.dumps(list_history(), ensure_ascii=False))
        elif self.path == "/favicon.ico":
            self._send(204, b"")
        elif self.path == "/api/render-status":
            with render_lock:
                snap = {k: v for k, v in render_status.items() if k != "_proc"}
            self._send(200, json.dumps(snap))
        elif self.path == "/api/root":
            # 让一键启动脚本检测当前实例的 root 是否匹配，不匹配则自动重启
            self._send(200, json.dumps({"root": ALLOWED_ROOT, "default": DEFAULT_FILE}, ensure_ascii=False))
        elif self.path == "/api/feishu-auth":
            self._send(200, json.dumps(run_feishu_auth(), ensure_ascii=False))
        elif self.path == "/api/download-pptx":
            self._serve_pptx()
        else:
            # fallback: 静态文件（slide 引用的 ./guizang-fde-red/*.png 等配图），仅放行媒体/样式类扩展名
            rel = self.path.lstrip("/").split("?")[0]
            if rel.startswith("..") or "\x00" in rel or ".fde_history" in rel.split("/"):
                self._send(403, "forbidden")
                return
            ext = os.path.splitext(rel)[1].lower()
            if ext not in (".png", ".jpg", ".jpeg", ".webp", ".svg", ".gif",
                           ".css", ".js", ".html", ".woff2", ".woff", ".ttf"):
                self._send(404, "not found: " + rel)
                return
            fp = os.path.join(ROOT, rel)
            if os.path.isfile(fp):
                ctype = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                         ".webp": "image/webp", ".svg": "image/svg+xml", ".gif": "image/gif",
                         ".css": "text/css", ".js": "application/javascript",
                         ".html": "text/html", ".woff2": "font/woff2",
                         ".woff": "font/woff", ".ttf": "font/ttf"}.get(ext, "application/octet-stream")
                with open(fp, "rb") as f:
                    self._send(200, f.read(), ctype)
            else:
                self._send(404, "not found: " + rel)

    def do_POST(self):
        length = int(self.headers.get("Content-Length", 0))
        raw = self.rfile.read(length) if length else b"{}"
        try:
            data = json.loads(raw.decode("utf-8")) if raw.strip() else {}
        except Exception:
            data = {}
        if self.path == "/api/save":
            rel = data.get("file", DEFAULT_FILE)
            html = data.get("html", "")
            fp = resolve_file(rel)
            if not fp:
                self._send(400, json.dumps({"error": "非法文件或未打开 deck"}))
                return
            try:
                tmp = fp + ".tmp"
                with open(tmp, "w", encoding="utf-8") as f:
                    f.write(html)
                os.replace(tmp, fp)
                save_snapshot(rel, html)   # 每次保存自动留历史快照，方便手动回滚
                self._send(200, json.dumps({"ok": True, "bytes": len(html)}))
            except Exception as e:
                self._send(500, json.dumps({"error": str(e)}))
        elif self.path == "/api/render":
            rel = data.get("file", DEFAULT_FILE)
            fp = resolve_file(rel)
            if not fp:
                self._send(400, json.dumps({"ok": False, "msg": "非法文件或未打开 deck"}))
                return
            if render_status["running"]:
                self._send(200, json.dumps({"ok": False, "msg": "渲染进行中，请稍候"}))
                return
            t = threading.Thread(target=run_render, args=(fp,), daemon=True)
            t.start()
            self._send(200, json.dumps({"ok": True, "msg": "已启动后台渲染"}))
        elif self.path == "/api/cancel-render":
            with render_lock:
                proc = render_status.get("_proc")
                if proc is not None:
                    try:
                        proc.kill()
                    except Exception:
                        pass
                render_status["cancelled"] = True
            self._send(200, json.dumps({"ok": True}))
        elif self.path == "/api/export-pptx":
            rel = data.get("file", DEFAULT_FILE)
            mode = data.get("mode", "editable")
            fp = resolve_file(rel)
            if not fp:
                self._send(400, json.dumps({"ok": False, "msg": "非法文件或未打开 deck"}))
                return
            if render_status["running"]:
                self._send(200, json.dumps({"ok": False, "msg": "任务进行中，请稍候"}))
                return
            t = threading.Thread(target=run_export, args=(fp, mode), daemon=True)
            t.start()
            self._send(200, json.dumps({"ok": True, "msg": "已启动导出 PPT"}))
        elif self.path == "/api/rollback":
            rel = data.get("file", DEFAULT_FILE)
            snap = data.get("snap", "")
            fp = resolve_file(rel)
            if not fp or not safe_snap(snap):
                self._send(400, json.dumps({"error": "非法参数"}))
                return
            snap_fp = os.path.join(HISTORY_BASE, deck_key(ACTIVE_FILE), snap)
            sb = os.path.normpath(snap_fp)
            hb = os.path.normpath(HISTORY_BASE)
            if not sb.startswith(hb + os.sep) or not os.path.isfile(sb):
                self._send(404, json.dumps({"error": "快照不存在"}))
                return
            try:
                # 回滚前先把当前源留一份快照，保证永可再回退
                cur = open(fp, "r", encoding="utf-8").read()
                save_snapshot(rel, cur)
            except Exception:
                pass
            try:
                content = open(sb, "r", encoding="utf-8").read()
                tmp = fp + ".tmp"
                with open(tmp, "w", encoding="utf-8") as f:
                    f.write(content)
                os.replace(tmp, fp)
                self._send(200, json.dumps({"ok": True, "html": content}))
            except Exception as e:
                self._send(500, json.dumps({"error": str(e)}))
        elif self.path == "/api/upload-feishu":
            kind = data.get("kind", "html")
            res = run_upload_feishu(kind)
            self._send(200, json.dumps(res, ensure_ascii=False))
        elif self.path == "/api/recent-remove":
            p = data.get("path", "")
            remove_recent(p)
            self._send(200, json.dumps({"ok": True}))
        elif self.path == "/api/download-pptx":
            self._serve_pptx()
        else:
            self._send(404, "not found")

    def _serve_pptx(self):
        with render_lock:
            name = render_status.get("pptx_name", "")
        if not name:
            self._send(404, "not ready")
            return
        fp = os.path.join(RENDER_BASE, deck_key(ACTIVE_FILE), name) if ACTIVE_FILE else None
        if not fp or not os.path.isfile(fp):
            self._send(404, "file missing")
            return
        with open(fp, "rb") as f:
            data = f.read()
        self.send_response(200)
        self.send_header("Content-Type",
                         "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        # HTTP 头只能 latin-1，中文文件名用 RFC 5987 (filename*=UTF-8'') 编码，并附 ASCII 兜底名
        self.send_header("Content-Disposition",
                         f'attachment; filename="fde-pptx.pptx"; filename*=UTF-8\'\'{quote(name)}')
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def log_message(self, *a):
        pass


if __name__ == "__main__":
    socketserver.TCPServer.allow_reuse_address = True
    with socketserver.ThreadingTCPServer(("127.0.0.1", PORT), Handler) as httpd:
        print(f"帆哥 PPT 编辑器 · http://localhost:{PORT}")
        print(f"  编辑器目录 = {EDITOR_HOME}")
        print(f"  打开 deck 后，root 会自动切到该 deck 所在目录（原地编辑，保存写回原文件）")
        print(f"  历史快照 / 渲染产物存于编辑器目录，绝不写入你的项目文件夹")
        try:
            httpd.serve_forever()
        except KeyboardInterrupt:
            print("\nbye.")
